"""J.7 / J.8 PowerPoint extraction."""

from __future__ import annotations

import asyncio
import io
import logging

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.group import GroupShape

from trimcp.extractors.core import ExtractionResult, Section, empty_skipped
from trimcp.extractors.libreoffice import libreoffice_convert
from trimcp.extractors.ocr import ocr_image_bytes

log = logging.getLogger(__name__)

EMU_LARGE_IMAGE = 2_000_000  # J.7 heuristic (~200x200 pt)


def _slide_hidden(slide) -> bool:
    try:
        show = slide._element.get("show")
        return show in ("0", "false", "False")
    except Exception:
        return False


async def _shape_parts(shape, warnings: list[str]) -> list[str]:
    parts: list[str] = []
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            if isinstance(shape, GroupShape):
                for child in shape.shapes:
                    parts.extend(await _shape_parts(child, warnings))
            return parts
    except Exception as e:
        warnings.append(f"group_shape_skip:{e}")
        return parts

    try:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                parts.append(txt)
    except Exception as e:
        warnings.append(f"text_frame_skip:{e}")

    try:
        if getattr(shape, "has_table", False) and shape.has_table:
            rows = [
                "| "
                + " | ".join(cell.text.replace("|", "\\|") for cell in row.cells)
                + " |"
                for row in shape.table.rows
            ]
            parts.append("\n".join(rows))
    except Exception as e:
        warnings.append(f"chart_or_table_skip:{e}")

    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if shape.width and int(shape.width) > EMU_LARGE_IMAGE:
                blob = await asyncio.to_thread(lambda: shape.image.blob)
                txt, ow = await ocr_image_bytes(blob)
                warnings.extend(ow)
                if txt:
                    parts.append(f"[image text: {txt}]")
    except Exception as e:
        warnings.append(f"picture_ocr_skip:{e}")

    return parts


async def extract_pptx(blob: bytes) -> ExtractionResult:
    warnings: list[str] = []

    def _load():
        return Presentation(io.BytesIO(blob))

    try:
        prs = await asyncio.to_thread(_load)
    except Exception as e:
        log.warning("python-pptx failed: %s", e)
        return empty_skipped("python-pptx", "corrupt", warnings=[str(e)])

    sections: list[Section] = []
    order = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        if _slide_hidden(slide):
            warnings.append(f"Skipped hidden slide: {slide_num}")
            continue
        try:
            title_shape = slide.shapes.title
        except AttributeError:
            title_shape = None
        slide_text_parts: list[str] = []
        try:
            if title_shape is not None and title_shape.has_text_frame:
                slide_text_parts.append(f"# {title_shape.text_frame.text}")
        except Exception as e:
            warnings.append(f"slide_title_skip:{e}")

        try:
            others = [s for s in slide.shapes if s is not title_shape]
            sorted_shapes = sorted(
                others, key=lambda s: (int(s.top or 0), int(s.left or 0))
            )
        except Exception:
            sorted_shapes = list(slide.shapes)

        for shape in sorted_shapes:
            try:
                slide_text_parts.extend(await _shape_parts(shape, warnings))
            except Exception as e:
                warnings.append(f"shape_walk_skip:{e}")

        body = "\n\n".join(p for p in slide_text_parts if p)
        sections.append(
            Section(
                text=body,
                structure_path=f"Slide {slide_num}",
                section_type="slide",
                order=order,
            )
        )
        order += 1

        try:
            if (
                slide.has_notes_slide
                and slide.notes_slide.notes_text_frame.text.strip()
            ):
                sections.append(
                    Section(
                        text=slide.notes_slide.notes_text_frame.text,
                        structure_path=f"Slide {slide_num} — Speaker Notes",
                        section_type="note",
                        order=order,
                    )
                )
                order += 1
        except Exception as e:
            warnings.append(f"notes_skip:{e}")

    meta: dict = {}
    try:
        core = prs.core_properties
        meta = {
            "title": core.title,
            "author": core.author,
            "created": core.created.isoformat() if core.created else None,
            "modified": core.modified.isoformat() if core.modified else None,
        }
    except Exception:
        pass

    full = "\n\n".join(s.text for s in sections)
    return ExtractionResult(
        method="python-pptx",
        text=full,
        sections=sections,
        metadata=meta,
        warnings=warnings,
    )


async def extract_ppt(blob: bytes) -> ExtractionResult:
    converted = await asyncio.to_thread(libreoffice_convert, blob, ".ppt", ".pptx")
    if not converted:
        return empty_skipped(
            "libreoffice", "conversion_failed", warnings=["ppt conversion failed"]
        )
    res = await extract_pptx(converted)
    res.method = "libreoffice→python-pptx"
    res.warnings.insert(0, "Converted from legacy .ppt via LibreOffice")
    return res
